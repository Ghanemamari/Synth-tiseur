import streamlit as st
import os, io, re, zipfile, tempfile
import fitz  # PyMuPDF
import docx  # python-docx
import pptx  # python-pptx
import pandas as pd
import pytesseract
from PIL import Image
import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
from unidecode import unidecode
from openai import OpenAI
import openai
from fpdf import FPDF  # Pour générer le PDF

# --- CLIENT API NVIDIA ---
def get_client():
    return OpenAI(
        base_url="https://integrate.api.nvidia.com/v1",
        api_key="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"
    )

# La fonction call_llm n'est plus utilisée pour l'affichage progressif.
def call_llm(prompt):
    client = get_client()
    completion = client.chat.completions.create(
        model="mistralai/mistral-small-24b-instruct",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.2,
        top_p=0.7,
        max_tokens=1024,
        stream=True
    )
    result = ""
    for chunk in completion:
        if chunk.choices[0].delta.content is not None:
            result += chunk.choices[0].delta.content
    return result.strip()

# --- FONCTIONS D'EXTRACTION DE TEXTE ---
def extract_text_from_pdf(file_obj):
    text = ""
    try:
        doc = fitz.open(stream=file_obj.read(), filetype="pdf")
        for page in doc:
            text += page.get_text() + "\n"
    except Exception as e:
        st.error(f"Erreur lors de l'extraction du PDF: {e}")
    return text

def extract_text_from_docx(file_obj):
    text = ""
    try:
        doc = docx.Document(file_obj)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        st.error(f"Erreur lors de l'extraction du DOCX: {e}")
    return text

def extract_text_from_pptx(file_obj):
    text = ""
    try:
        presentation = pptx.Presentation(file_obj)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        st.error(f"Erreur lors de l'extraction du PPTX: {e}")
    return text

def extract_text_from_xlsx(file_obj):
    text = ""
    try:
        df = pd.read_excel(file_obj)
        text = df.to_string()
    except Exception as e:
        st.error(f"Erreur lors de l'extraction du XLSX: {e}")
    return text

def extract_text_from_txt(file_obj):
    try:
        text = file_obj.read().decode("utf-8")
    except Exception as e:
        st.error(f"Erreur lors de l'extraction du TXT: {e}")
        text = ""
    return text

def extract_text_from_image(file_obj):
    text = ""
    try:
        image = Image.open(file_obj)
        text = pytesseract.image_to_string(image)
    except Exception as e:
        st.error(f"Erreur lors de l'extraction de l'image: {e}")
    return text

def extract_text_from_zip(file_obj):
    full_text = ""
    try:
        with zipfile.ZipFile(file_obj) as z:
            for name in z.namelist():
                ext = name.split('.')[-1].lower()
                if ext in ["pdf", "docx", "pptx", "xlsx", "txt", "jpg", "jpeg", "png"]:
                    with z.open(name) as f:
                        full_text += process_file(f, name) + "\n"
    except Exception as e:
        st.error(f"Erreur lors de l'extraction du zip: {e}")
    return full_text

def process_file(file_obj, file_name):
    ext = file_name.split('.')[-1].lower()
    if ext == "pdf":
        return extract_text_from_pdf(file_obj)
    elif ext == "docx":
        return extract_text_from_docx(file_obj)
    elif ext == "pptx":
        return extract_text_from_pptx(file_obj)
    elif ext == "xlsx":
        return extract_text_from_xlsx(file_obj)
    elif ext == "txt":
        return extract_text_from_txt(file_obj)
    elif ext in ["jpg", "jpeg", "png"]:
        return extract_text_from_image(file_obj)
    elif ext == "zip":
        return extract_text_from_zip(file_obj)
    else:
        return ""

# --- FONCTIONS DE TRAITEMENT DU TEXTE ---
def clean_text(text):
    text = text.lower()
    text = unidecode(text)
    text = re.sub(r"[^a-z0-9\s.,;:?!'-]", " ", text)
    text = re.sub(r'\s+', ' ', text)
    return text

def split_text_into_chunks(text, chunk_size=500, overlap=50):
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        if end >= len(words):
            break
        start = end - overlap
    return chunks

# --- FONCTIONS D'EMBEDDING ET D'INDEXATION ---
def embed_chunks(chunks, model):
    embeddings = model.encode(chunks, convert_to_numpy=True)
    return embeddings

def build_faiss_index(embeddings):
    d = embeddings.shape[1]
    index = faiss.IndexFlatL2(d)
    index.add(embeddings)
    return index

def search_best_chunks(index, query_embedding, k=8):
    distances, indices = index.search(np.array([query_embedding]), k)
    return indices[0]

# --- FONCTION DE GENERATION DE PROMPT ---
def generate_prompt(context, language="francais"):
    if language == "francais":
        prompt = (
            "Tu es un assistant specialise en recherche scientifique. Resume le document suivant (article scientifique ou these) en suivant rigoureusement la structure ci-dessous :\n\n"
            "Contexte et objectifs : Quelle est la problematique aborde et quels sont les objectifs de l’etude ?\n"
            "Verrous : Quels sont les principaux verrous scientifiques, techniques ou methodologiques identifies ?\n"
            "Approche globale : Quelle est la strategie ou la methodologie adoptee ?\n"
            "Principales iterations et performances associees : Quelles experiences ont ete menees et avec quels resultats ?\n"
            "Perspectives : Quelles sont les pistes d’amelioration ou perspectives futures ?\n\n"
            "Voici le contenu pertinent extrait du document :\n"
            "[context]\n\n"
            "Redige maintenant le resume structure."
        )
    else:
        prompt = (
            "You are a scientific assistant. Summarize the following document strictly following this structure:\n\n"
            "1. **Context and Objectives**\n"
            "2. **Challenges**\n"
            "3. **Global Approach**\n"
            "4. **Main Iterations and Associated Performance**\n"
            "5. **Perspectives**\n\n"
            "Here is the relevant content extracted from the document:\n"
            "[context]\n\n"
            "Now, write the structured summary."
        )
    return prompt.replace("[context]", context)

# --- FONCTION DE GENERATION DU PDF ---
def generate_pdf(text):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, text)
    pdf_bytes = pdf.output(dest="S").encode("latin1")
    return pdf_bytes

# --- INTERFACE STREAMLIT ---
st.title("Résumé automatique de documents scientifiques")
st.markdown("Cette application vous permet de générer automatiquement un résumé structuré de documents scientifiques.")

# --- Sélection de la langue placée avant l'upload ---
language = st.selectbox("Langue du résumé", ["Français", "Anglais"])
lang_choice = "francais" if language == "Français" else "english"

# --- Partie d'upload ---
option = st.radio("Que veux-tu résumer ?", ["Fichier unique", "Dossier (ou .zip)"])
allowed_types = ["pdf", "docx", "pptx", "xlsx", "txt", "jpg", "jpeg", "png", "zip"]
uploaded_files = st.file_uploader(
    "📎 Charge ton fichier ou dossier :",
    type=allowed_types,
    accept_multiple_files=(option == "Dossier (ou .zip)")
)

if st.button("Générer le résumé"):
    if uploaded_files:
        st.info("Extraction et traitement des documents en cours...")
        # On convertit uploaded_files en liste pour un traitement homogène
        if option == "Fichier unique":
            files_list = [uploaded_files] if not isinstance(uploaded_files, list) else uploaded_files
        else:
            files_list = uploaded_files

        all_texts = []
        for file_obj in files_list:
            file_name = file_obj.name
            if file_name.split('.')[-1].lower() == "zip":
                text = extract_text_from_zip(file_obj)
            else:
                text = process_file(file_obj, file_name)
            all_texts.append(text)
        
        full_text = "\n".join(all_texts)
        
        if not full_text.strip():
            st.error("Aucun texte n'a pu être extrait des documents fournis.")
        else:
            cleaned_text = clean_text(full_text)
            st.text_area("Texte nettoyé (aperçu)", cleaned_text[:500] + " ...", height=150)
            
            chunks = split_text_into_chunks(cleaned_text, chunk_size=500, overlap=50)
            st.write(f"{len(chunks)} chunks générés.")
            
            with st.spinner("Calcul des embeddings..."):
                model = SentenceTransformer("all-MiniLM-L6-v2")
                embeddings = embed_chunks(chunks, model)
            
            index = build_faiss_index(np.array(embeddings))
            
            query = "document scientifique avec objectifs, verrous, méthodologie, résultats, perspectives"
            query_embedding = model.encode(query, convert_to_numpy=True)
            best_indices = search_best_chunks(index, query_embedding, k=8)
            
            selected_chunks = [chunks[i] for i in best_indices if i < len(chunks)]
            context = "\n\n".join(selected_chunks)
            
            prompt = generate_prompt(context, language=lang_choice)
            
            st.markdown("**Appel à l'API LLM en cours...**")
            # Affichage progressif du résumé, en mode streaming
            with st.spinner("Génération du résumé..."):
                try:
                    client = get_client()
                    response = client.chat.completions.create(
                        model="mistralai/mistral-small-24b-instruct",
                        messages=[{"role": "user", "content": prompt}],
                        temperature=0.2,
                        top_p=0.7,
                        max_tokens=1024,
                        stream=True
                    )
                    st.subheader("🧾 Résumé généré :")
                    summary = ""
                    summary_container = st.empty()
                    for chunk in response:
                        if chunk.choices[0].delta.content:
                            summary += chunk.choices[0].delta.content
                            summary_container.write(summary)
                    st.success("Résumé terminé.")
                    st.download_button("Télécharger le résumé en PDF",
                                       data=generate_pdf(summary),
                                       file_name="resume.pdf",
                                       mime="application/pdf")
                except Exception as e:
                    st.error(f"Erreur lors de l'appel au LLM : {e}")
    else:
        st.warning("Veuillez uploader un ou plusieurs fichiers selon le mode choisi.")
